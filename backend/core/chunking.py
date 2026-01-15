# backend/core/chunking.py

from dataclasses import dataclass
from typing import List, Dict, Any, Tuple
from pathlib import Path
import uuid, re, tiktoken, os, logging

logger = logging.getLogger(__name__)

# Check Docling Availability
try:
    from docling.document_converter import DocumentConverter
    from docling.datamodel.base_models import InputFormat
    DOCLING_AVAILABLE = True
    logger.info(" Docling khả dụng cho xử lý PDF/DOCX")
except ImportError:
    DOCLING_AVAILABLE = False
    logger.error(" Docling không khả dụng! Xử lý PDF/DOCX sẽ thất bại.")

# Config
@dataclass
class ChunkConfig:
    max_tokens: int = 300
    overlap_tokens: int = 50
    lookback_ratio: float = 0.2
    table_rows_per_chunk: int = 50
    join_separator: str = "\n"

DocChunkConfig = ChunkConfig

# Constants
_SENTENCE_BREAK = re.compile(r"(?s)(.*?)([\\.!?…]|(?:\n{2,})|(?:\r?\n- )|(?:\r?\n• ))\s+$")

class DoclingExtractor:
    """Trích xuất văn bản từ PDF/DOCX bằng Docling"""
    
    _instance = None
    
    def __new__(cls):
        if cls._instance is None:
            cls._instance = super().__new__(cls)
            cls._instance._initialized = False
        return cls._instance
    
    def __init__(self):
        if self._initialized:
            return
        
        if not DOCLING_AVAILABLE:
            raise ImportError(" Docling chưa được cài đặt!")
        
        try:
            self.converter = DocumentConverter()
            logger.info(" Đã khởi tạo Docling (sử dụng API mới nhất)")
        except TypeError:
            try:
                from docling.datamodel.pipeline_options import PdfPipelineOptions
                device = os.getenv('DOCLING_DEVICE', 'cpu')
                do_ocr = os.getenv('DOCLING_OCR', 'false').lower() == 'true'
                
                pipeline_options = PdfPipelineOptions()
                pipeline_options.do_ocr = do_ocr
                pipeline_options.do_table_structure = True
                
                self.converter = DocumentConverter(
                    allowed_formats=[InputFormat.PDF, InputFormat.DOCX],
                    pipeline_options=pipeline_options
                )
                logger.info(f" Đã khởi tạo Docling ( thiết bị={device}, ocr={do_ocr})")
            except Exception as e:
                logger.error(f" Khởi tạo Docling thất bại: {e}")
                raise RuntimeError(f"Failed to initialize Docling: {e}")
        
        self._initialized = True
    
    def extract(self, filepath: str) -> Tuple[str, dict]:
        """Trích xuất văn bản và metadata từ file"""
        try:
            logger.info(f" Docling đang trích xuất: {Path(filepath).name}")
            result = self.converter.convert(filepath)
            metadata = self._extract_metadata(result)
            
            # Thử các phương thức xuất: markdown -> text -> custom
            for method, func in [
                ("markdown", lambda: result.document.export_to_markdown()),
                ("text", lambda: result.document.export_to_text()),
                ("custom", lambda: self._custom_export(result))
            ]:
                try:
                    text = func()
                    if text and text.strip():
                        metadata['export_method'] = method
                        logger.info(f" Xuất {method}: {len(text)} ký tự, {len(text.split())} từ")
                        return text, metadata
                except Exception as e:
                    if method == "custom":
                        raise
                    logger.warning(f" Xuất {method} thất bại: {e}")
            
            return "", metadata
        except Exception as e:
            logger.error(f" Trích xuất Docling thất bại cho {Path(filepath).name}: {e}")
            raise RuntimeError(f"Docling extraction failed: {e}")
    
    def _extract_metadata(self, result) -> dict:
        """Trích xuất metadata từ tài liệu"""
        metadata = {}
        try:
            doc = result.document
            
            # Ánh xạ tên thuộc tính sang key metadata
            attr_map = {
                'title': ['title', 'name'],
                'author': ['author', 'authors'],
                'created_date': ['created_date', 'creation_date'],
                'modified_date': ['modified_date'],
                'keywords': ['keywords'],
                'subject': ['subject'],
                'page_count': ['page_count', 'num_pages']
            }
            
            for key, attrs in attr_map.items():
                for attr in attrs:
                    if hasattr(doc, attr):
                        value = getattr(doc, attr)
                        if value:
                            if isinstance(value, list):
                                metadata[key] = ', '.join(value)
                            else:
                                metadata[key] = str(value)
                            break
            
            # Gộp metadata bổ sung nếu có
            if hasattr(doc, 'metadata') and isinstance(doc.metadata, dict):
                for k, v in doc.metadata.items():
                    if k not in metadata and v:
                        metadata[k] = str(v)
            
            if metadata:
                logger.info(f" Đã trích xuất {len(metadata)} trường metadata: {list(metadata.keys())}")
        except Exception as e:
            logger.warning(f" Không thể trích xuất metadata: {e}")
        
        return metadata
    
    def _custom_export(self, result) -> str:
        """Xuất văn bản có cấu trúc tùy chỉnh"""
        lines = []
        
        # Element type to markdown mapping
        element_map = {
            "title": lambda t: f"# {t}",
            "section_header": lambda t: f"## {t}",
            "subtitle": lambda t: f"### {t}",
            **{f"heading_{i}": lambda t, i=i: f"{'#' * i} {t}" for i in range(1, 7)},
            **{f"h{i}": lambda t, i=i: f"{'#' * i} {t}" for i in range(1, 7)},
            "list_item": lambda t: f"- {t}",
            "numbered_list": lambda t: f"1. {t}",
            "ordered_list": lambda t: f"1. {t}",
            "bullet_list": lambda t: f"- {t}",
            "unordered_list": lambda t: f"- {t}",
            "code_block": lambda t: f"```\n{t}\n```",
            "code": lambda t: f"```\n{t}\n```",
            "pre": lambda t: f"```\n{t}\n```",
            "formula": lambda t: f"$$\n{t}\n$$",
            "equation": lambda t: f"$$\n{t}\n$$",
            "math": lambda t: f"$$\n{t}\n$$",
            "quote": lambda t: "\n".join(f"> {line}" for line in t.split("\n")),
            "block_quote": lambda t: "\n".join(f"> {line}" for line in t.split("\n")),
            "blockquote": lambda t: "\n".join(f"> {line}" for line in t.split("\n")),
            "figure": lambda t: f"![Figure: {t}]",
            "picture": lambda t: f"![Figure: {t}]",
            "image": lambda t: f"![Figure: {t}]",
            "page_header": lambda t: f"*[Header: {t}]*",
            "header": lambda t: f"*[Header: {t}]*",
            "page_footer": lambda t: f"*[Footer: {t}]*",
            "footer": lambda t: f"*[Footer: {t}]*",
            "reference": lambda t: f"[^{t}]",
            "citation": lambda t: f"[^{t}]",
            "bibliography": lambda t: f"[^{t}]",
            "text_box": lambda t: f" **Note:** {t}",
            "callout": lambda t: f" **Note:** {t}",
            "note": lambda t: f"**Note:** {t}",
            "caption": lambda t: f"*{t}*",
            "footnote": lambda t: f"[^{t}]",
            "paragraph": lambda t: t
        }
        
        try:
            if hasattr(result.document, 'iterate_items'):
                for element in result.document.iterate_items():
                    label = getattr(element, 'label', 'unknown')
                    text = getattr(element, 'text', '').strip()
                    
                    if not text:
                        continue
                    
                    # Handle tables specially
                    if label == "table":
                        if hasattr(element, 'export_to_dataframe'):
                            try:
                                df = element.export_to_dataframe()
                                lines.append(df.to_markdown(index=False))
                                logger.debug(f" Bảng đã xuất dạng markdown ({len(df)} dòng)")
                            except Exception as e:
                                logger.warning(f" Xuất bảng thất bại: {e}, sử dụng văn bản thô")
                                lines.append(f"```\n{text}\n```")
                        else:
                            lines.append(f"```\n{text}\n```")
                    # Use mapping for other elements
                    elif label in element_map:
                        lines.append(element_map[label](text))
                    else:
                        logger.warning(f" Loại phần tử không xác định '{label}', giữ nguyên dạng văn bản")
                        lines.append(text)
                    
                    lines.append("")  # Add spacing
            
            if lines:
                return "\n".join(lines)
            
            # Fallback methods
            for attr in ['text', 'to_dict']:
                if hasattr(result.document, attr):
                    logger.warning(f" Sử dụng document.{attr} dự phòng")
                    if attr == 'text':
                        return result.document.text
                    else:
                        return str(result.document.to_dict().get('text', ''))
            
            logger.error(" Tất cả phương thức xuất đều thất bại")
            return ""
        except Exception as e:
            logger.error(f" Xuất tùy chỉnh thất bại: {e}")
            raise RuntimeError(f"Cannot extract text from document: {e}")


def _with_breadcrumb(section: str, content: str, part_idx: int, part_total: int) -> str:
    """Thêm breadcrumb section vào chunk"""
    suffix = f" — part {part_idx}/{part_total}" if part_total > 1 else ""
    return f"**[SECTION] {section}{suffix}**\n\n{content}"


def _soft_split(text: str, enc, max_size: int, overlap: int, lookback_ratio: float = 0.2) -> List[str]:
    """Chia văn bản thông minh theo ranh giới câu"""
    ids = enc.encode(text)
    n = len(ids)
    if n <= max_size:
        return [text]

    out, start = [], 0
    while start < n:
        end = min(start + max_size, n)
        window_ids = ids[start:end]
        window_text = enc.decode(window_ids)

        if end < n:
            lb_chars = max(10, int(len(window_text) * (1 - lookback_ratio)))
            tail = window_text[lb_chars:]
            m = _SENTENCE_BREAK.search(tail)
            if m:
                cut_char = lb_chars + m.end()
                window_ids = enc.encode(window_text[:cut_char])

        out.append(enc.decode(window_ids).rstrip())
        if start + len(window_ids) >= n:
            break
        start += len(window_ids) - overlap
    return out


def _split_table_text(header: str, rows: List[str], prefix: str, enc, max_tokens: int, overlap: int) -> List[str]:
    """Chia bảng với header cố định"""
    sticky = header + "\n"
    chunks, cur_rows = [], []

    def emit():
        if not cur_rows:
            return
        body = "".join([r + "\n" for r in cur_rows])
        content = f"{prefix}\n\n{sticky}{body}".strip() if prefix else f"{sticky}{body}".strip()
        if len(enc.encode(content)) > max_tokens:
            chunks.extend(_soft_split(content, enc, max_tokens, overlap))
        else:
            chunks.append(content)

    for r in rows:
        candidate = cur_rows + [r]
        body = "".join([x + "\n" for x in candidate])
        content = f"{prefix}\n\n{sticky}{body}".strip() if prefix else f"{sticky}{body}".strip()
        if len(enc.encode(content)) <= max_tokens:
            cur_rows.append(r)
        else:
            emit()
            cur_rows = [r]
    emit()
    if not chunks:
        chunks.append(f"{prefix}\n\n{sticky}".strip() if prefix else sticky.strip())
    return chunks


def _split_table_markdown(text: str, enc, max_tokens: int, overlap: int) -> List[str]:
    """Chia bảng markdown thông minh"""
    lines = text.splitlines()
    block_start, block_end = None, None
    
    for i in range(len(lines) - 1):
        s1, s2 = lines[i].strip(), lines[i + 1].strip()
        if s1.startswith("|") and s1.endswith("|") and all(c in "-: " for c in s2.replace("|", "")):
            block_start, block_end = i, i + 2
            while block_end < len(lines) and lines[block_end].lstrip().startswith("|"):
                block_end += 1
            break
    
    if block_start is None:
        return _soft_split(text, enc, max_tokens, overlap)

    prefix = "\n".join(lines[:block_start]).strip()
    table_lines = lines[block_start:block_end]
    header, separator = table_lines[:2]
    rows = table_lines[2:]
    suffix = "\n".join(lines[block_end:]).strip()
    
    # Process table
    table_chunks = _split_table_text(header + "\n" + separator, rows, prefix, enc, max_tokens, overlap)
    
    # Append suffix
    if suffix:
        if table_chunks:
            last_chunk = table_chunks[-1]
            combined = last_chunk + "\n\n" + suffix
            if len(enc.encode(combined)) <= max_tokens:
                table_chunks[-1] = combined
            else:
                table_chunks.extend(_soft_split(suffix, enc, max_tokens, overlap))
        else:
            table_chunks = _soft_split(suffix, enc, max_tokens, overlap)
    
    return table_chunks


class Chunker:
    """Chia văn bản theo token với theo dõi cấu trúc heading"""
    
    def __init__(self, config: ChunkConfig = None):
        if config is None:
            try:
                from backend.config import Config
                config = ChunkConfig(
                    max_tokens=Config.DEFAULT_CHUNK_SIZE,
                    overlap_tokens=Config.DEFAULT_CHUNK_OVERLAP
                )
            except:
                config = ChunkConfig()
        
        self.config = config
        self.enc = tiktoken.encoding_for_model("gpt-4o-mini")
        self.heading_stack = []

    def count_tokens(self, text: str) -> int:
        """Đếm số token trong văn bản"""
        return len(self.enc.encode(text))
    
    def _parse_heading(self, line: str) -> Tuple[int, str]:
        """Phân tích heading markdown, trả về (level, title) hoặc (0, None)"""
        line = line.strip()
        if line.startswith('#'):
            level = len(line) - len(line.lstrip('#'))
            if level <= 6 and level < len(line) and line[level] == ' ':
                return (level, line[level:].strip())
        return (0, None)
    
    def _update_heading_stack(self, level: int, title: str):
        """Cập nhật stack cấu trúc heading"""
        self.heading_stack = [h for h in self.heading_stack if h[0] < level]
        self.heading_stack.append((level, title))
    
    def _get_current_section(self) -> str:
        """Lấy đường dẫn section hiện tại từ heading stack"""
        return " > ".join(h[1] for h in self.heading_stack) if self.heading_stack else "ROOT"

    def chunk_text(self, text: str, filepath: str, section: str = "ROOT") -> List[Dict]:
        """Chia văn bản theo token, xử lý bảng và theo dõi cấu trúc heading"""
        input_chars = len(text)
        input_tokens = self.count_tokens(text)
        logger.info(f" Input: {input_chars} ký tự, {len(text.split())} từ, {input_tokens} tokens")
        
        lines = text.split('\n')
        chunks, buf, buf_tokens, order = [], [], 0, 0
        self.heading_stack = []
        tiny_chunks_count = merged_chunks_count = 0

        def process_buffer():
            nonlocal order, tiny_chunks_count, merged_chunks_count
            if not buf:
                return []
            
            combined = "\n".join(buf).strip()
            pieces = _split_table_markdown(combined, self.enc, self.config.max_tokens, self.config.overlap_tokens)
            current_section = self._get_current_section()
            result = []
            
            for i, p in enumerate(pieces, 1):
                # Handle tiny chunks
                if len(p.strip()) < 5:
                    tiny_chunks_count += 1
                    logger.debug(f" Tiny chunk ({len(p)} chars): '{p[:30]}...'")
                    if chunks:
                        chunks[-1]['content'] += "\n" + p
                        chunks[-1]['tokens'] = self.count_tokens(chunks[-1]['content'])
                        merged_chunks_count += 1
                        logger.debug(f" Merged tiny chunk vào chunk trước ({chunks[-1]['tokens']} tokens)")
                        continue
                    else:
                        logger.debug(f" Giữ tiny chunk cho chunk tiếp theo")
                        return [p]
                
                result.append({
                    "chunk_id": str(uuid.uuid4()),
                    "content": _with_breadcrumb(current_section, p, i, len(pieces)),
                    "tokens": self.count_tokens(p),
                    "order": order,
                    "file_path": filepath,
                    "file_type": Path(filepath).suffix[1:].upper(),
                    "section": current_section
                })
                order += 1
            return result

        for line in lines:
            level, title = self._parse_heading(line)
            if level > 0:
                self._update_heading_stack(level, title)
            
            line_tokens = self.count_tokens(line)
            if buf_tokens + line_tokens > self.config.max_tokens and buf:
                chunks.extend(process_buffer())
                # Keep overlap
                overlap_lines = buf[-5:] if len(buf) > 5 else buf
                buf = overlap_lines + [line]
                buf_tokens = sum(self.count_tokens(l) for l in buf)
            else:
                buf.append(line)
                buf_tokens += line_tokens

        # Process remaining buffer
        chunks.extend(process_buffer())
        
        output_chars = sum(len(c['content']) for c in chunks)
        output_tokens = sum(c['tokens'] for c in chunks)
        
        logger.info(f" Output: {len(chunks)} chunks, {output_chars} ký tự, {output_tokens} tokens")
        logger.info(f" Retention: {output_chars}/{input_chars} chars ({output_chars/max(input_chars,1)*100:.1f}%), "
                   f"{output_tokens}/{input_tokens} tokens ({output_tokens/max(input_tokens,1)*100:.1f}%)")
        
        if tiny_chunks_count > 0:
            logger.info(f" Tiny chunks: {tiny_chunks_count} phát hiện, {merged_chunks_count} đã merge")
        
        retention_rate = output_chars / max(input_chars, 1)
        if retention_rate < 0.95:
            logger.warning(f" Mất {(1-retention_rate)*100:.1f}% nội dung trong quá trình chunking!")

        return chunks


# Hàm trích xuất file
def _read_file(filepath: str, encoding: str = 'utf-8') -> str:
    """Hàm đọc file chung với encoding dự phòng"""
    try:
        with open(filepath, 'r', encoding=encoding) as f:
            return f.read()
    except UnicodeDecodeError:
        if encoding != 'latin-1':
            return _read_file(filepath, 'latin-1')
        raise
    except Exception as e:
        logger.error(f" Lỗi đọc file: {e}")
        return ""


def _extract_markdown(filepath: str) -> str:
    return _read_file(filepath)


def _extract_html(filepath: str) -> str:
    try:
        from bs4 import BeautifulSoup
        html = _read_file(filepath)
        soup = BeautifulSoup(html, 'html.parser')
        for script in soup(["script", "style"]):
            script.decompose()
        return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        logger.error(f" Lỗi trích xuất HTML: {e}")
        return ""


def _extract_json(filepath: str) -> str:
    try:
        import json
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        logger.error(f" Lỗi trích xuất JSON: {e}")
        return ""


def _extract_xml(filepath: str) -> str:
    try:
        from bs4 import BeautifulSoup
        xml = _read_file(filepath)
        soup = BeautifulSoup(xml, 'xml')
        return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        logger.error(f" Lỗi trích xuất XML: {e}")
        return ""


def _extract_text(filepath: str) -> str:
    return _read_file(filepath)


def _extract_excel(filepath: str) -> str:
    try:
        import pandas as pd
        dfs = pd.read_excel(filepath, sheet_name=None)
        text_parts = []
        enc = tiktoken.encoding_for_model("gpt-4o-mini")
        
        for sheet_name, df in dfs.items():
            if df.empty:
                continue
            header = " | ".join(str(col) for col in df.columns)
            rows = [" | ".join(map(str, row)) for row in df.values]
            sheet_prefix = f"=== Sheet: {sheet_name} ==="
            text_parts.extend(_split_table_text(header, rows, sheet_prefix, enc, max_tokens=300, overlap=50))
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f" Lỗi trích xuất Excel: {e}")
        return ""


def _extract_csv(filepath: str) -> str:
    try:
        import pandas as pd
        df = pd.read_csv(filepath)
        if df.empty:
            return ""
        enc = tiktoken.encoding_for_model("gpt-4o-mini")
        header = " | ".join(str(col) for col in df.columns)
        rows = [" | ".join(map(str, row)) for row in df.values]
        return "\n\n".join(_split_table_text(header, rows, "", enc, max_tokens=300, overlap=50))
    except Exception as e:
        logger.error(f" Lỗi trích xuất CSV: {e}")
        return ""


def _extract_pdf_pypdf2(filepath: str) -> str:
    """Trích xuất PDF dự phòng bằng PyPDF2"""
    try:
        import PyPDF2
        text_parts = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"=== Page {page_num} ===\n{text}")
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f" Lỗi trích xuất PDF bằng PyPDF2: {e}")
        return ""


def _extract_pdf_pdfplumber(filepath: str) -> str:
    """Trích xuất PDF dự phòng bằng pdfplumber"""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(filepath) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    text_parts.append(f"=== Page {page_num} ===\n{text}")
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f" Lỗi trích xuất PDF bằng pdfplumber: {e}")
        return ""


def _extract_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text_parts = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = [shape.text.strip() for shape in slide.shapes 
                         if hasattr(shape, "text") and shape.text.strip()]
            if slide_text:
                text_parts.append(f"=== Slide {idx} ===\n" + "\n".join(slide_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f" Lỗi trích xuất PPTX: {e}")
        return ""


def extract_text_from_file(filepath: str) -> Tuple[str, dict]:
    """Trích xuất văn bản từ file với metadata
    
    Chiến lược:
    - PDF & DOCX: Thử Docling trước, dự phòng PyPDF2/pdfplumber nếu thất bại
    - Định dạng khác: Dùng các hàm trích xuất cũ
    
    Trả về:
        Tuple của (văn_bản, metadata_dict)
    """
    ext = Path(filepath).suffix.lower()
    metadata = {}
    
    # PDF/DOCX: Thử Docling với dự phòng
    if ext in ['.pdf', '.docx', '.doc']:
        # Try Docling first
        if DOCLING_AVAILABLE:
            try:
                logger.info(f" Đang thử Docling cho {ext}: {Path(filepath).name}")
                extractor = DoclingExtractor()
                text, metadata = extractor.extract(filepath)
                
                if text and text.strip():
                    logger.info(f" Docling đã trích xuất {len(text)} ký tự từ {Path(filepath).name}")
                    metadata['extraction_method'] = 'docling'
                    return text, metadata
                else:
                    logger.warning(f" Docling trả về văn bản rỗng, thử fallback...")
            except Exception as e:
                logger.warning(f" Docling thất bại: {e}, đang thử fallback...")
        
        # Dự phòng chỉ cho PDF
        if ext == '.pdf':
            # Try PyPDF2
            logger.info(f" Đang thử PyPDF2 cho PDF: {Path(filepath).name}")
            text = _extract_pdf_pypdf2(filepath)
            if text and text.strip():
                logger.info(f" PyPDF2 đã trích xuất {len(text)} ký tự")
                metadata['extraction_method'] = 'pypdf2'
                return text, metadata
            
            # Try pdfplumber
            logger.info(f" Đang thử pdfplumber cho PDF: {Path(filepath).name}")
            text = _extract_pdf_pdfplumber(filepath)
            if text and text.strip():
                logger.info(f" pdfplumber đã trích xuất {len(text)} ký tự")
                metadata['extraction_method'] = 'pdfplumber'
                return text, metadata
            
            raise RuntimeError(f" Tất cả phương pháp trích xuất PDF đều thất bại cho {Path(filepath).name}")
        
        # Với DOCX/DOC, Docling là bắt buộc
        raise RuntimeError(f" Docling là bắt buộc cho {ext} nhưng đã thất bại!")
    
    # Định dạng khác: Dùng các hàm trích xuất cũ
    extractors = {
        '.md': _extract_markdown, '.markdown': _extract_markdown,
        '.html': _extract_html, '.htm': _extract_html,
        '.json': _extract_json,
        '.xml': _extract_xml,
        '.xlsx': _extract_excel, '.xls': _extract_excel,
        '.csv': _extract_csv,
        '.pptx': _extract_pptx, '.ppt': _extract_pptx,
    }
    
    # Các file văn bản
    text_exts = ['.txt', '.py', '.js', '.java', '.cpp', '.c', '.h', '.hpp', 
                 '.css', '.scss', '.sql', '.sh', '.bash', '.yml', '.yaml',
                 '.toml', '.ini', '.cfg', '.conf', '.log', '.r', '.rb', 
                 '.php', '.go', '.rs', '.swift', '.kt', '.ts', '.jsx', '.tsx']
    
    for text_ext in text_exts:
        extractors[text_ext] = _extract_text
    
    if ext in extractors:
        return extractors[ext](filepath), metadata
    
    logger.warning(f" Loại file không được hỗ trợ: {ext}")
    return "", metadata


def process_document_to_chunks(filepath: str, config: ChunkConfig = None) -> Tuple[List[Dict], dict]:
    """Điểm vào chính cho xử lý tài liệu
    
    Tham số:
        filepath: Đường dẫn tới tài liệu
        config: Cấu hình chunking (dùng mặc định từ Config nếu None)
        
    Trả về:
        Tuple của (danh_sách_chunks, metadata_dict)
    """
    if config is None:
        try:
            from backend.config import Config
            config = ChunkConfig(
                max_tokens=Config.DEFAULT_CHUNK_SIZE,
                overlap_tokens=Config.DEFAULT_CHUNK_OVERLAP
            )
        except:
            config = ChunkConfig()
    
    logger.info(f" Đang xử lý: {filepath}")
    
    # Extract text
    text, metadata = extract_text_from_file(filepath)
    
    if not text.strip():
        logger.warning(f" Cảnh báo: Không trích xuất được văn bản từ {filepath}")
        return [], metadata
    
    logger.info(f" Đã trích xuất {len(text)} ký tự")
    
    # Chunk text
    chunker = Chunker(config)
    chunks = chunker.chunk_text(text, filepath)
    
    logger.info(f" Đã tạo {len(chunks)} chunks")
    return chunks, metadata


__all__ = [
    'ChunkConfig',
    'DocChunkConfig',
    'Chunker',
    'DoclingExtractor',
    'extract_text_from_file',
    'process_document_to_chunks',
    'DOCLING_AVAILABLE'
]